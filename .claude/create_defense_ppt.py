from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TEMPLATE = Path(r"D:/Browserdownload/产品介绍.pptx")
OUT = Path(r"D:/桌面/AICD-Platform2/AICD_Platform/AI_Creator_Platform_答辩PPT.pptx")

prs = Presentation(str(TEMPLATE))

TARGET_SLIDES = 22
while len(prs.slides) > TARGET_SLIDES:
    idx = len(prs.slides) - 1
    sld_id = prs.slides._sldIdLst[idx]
    rel_id = sld_id.rId
    prs.part.drop_rel(rel_id)
    prs.slides._sldIdLst.remove(sld_id)

while len(prs.slides) < TARGET_SLIDES:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    prs.slides.add_slide(blank_layout)

W = prs.slide_width
H = prs.slide_height
SLIDE_W = W / 914400
SLIDE_H = H / 914400

NAVY = RGBColor(14, 28, 56)
NAVY_2 = RGBColor(24, 44, 82)
BLUE = RGBColor(33, 115, 232)
CYAN = RGBColor(31, 205, 245)
ICE = RGBColor(232, 242, 255)
LIGHT = RGBColor(247, 250, 255)
TEXT = RGBColor(35, 47, 68)
MUTED = RGBColor(98, 114, 138)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(42, 180, 122)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(226, 78, 78)
GRAY = RGBColor(226, 232, 240)


def remove_all_shapes(slide):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


def fill_slide(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_textbox(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
                align=PP_ALIGN.LEFT, font="SimHei", valign="mid",
                line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multi(slide, lines, x, y, w, h, size=15, color=TEXT, header_color=None,
              font="SimHei", bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        if bullet:
            p.level = 0
        text = item if isinstance(item, str) else item[0]
        is_head = isinstance(item, tuple) and len(item) > 1 and item[1]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bool(is_head)
        run.font.color.rgb = header_color if is_head and header_color else color
    return box


def rect(slide, x, y, w, h, fill, line=None, radius=False, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, x, y, d, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def line(slide, x1, y1, x2, y2, color=BLUE, width=2):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def title(slide, t, sub=None, dark=False, idx=None):
    c = WHITE if dark else TEXT
    add_textbox(slide, t, 0.72, 0.42, 8.2, 0.48, size=24, color=c, bold=True)
    if sub:
        add_textbox(slide, sub, 0.74, 0.92, 8.4, 0.28, size=9.5, color=ICE if dark else MUTED)
    if idx:
        add_textbox(slide, f"{idx:02d}", 12.08, 0.42, 0.52, 0.28, size=10, color=ICE if dark else MUTED, align=PP_ALIGN.RIGHT)


def footer(slide, dark=False):
    color = RGBColor(170, 188, 215) if dark else RGBColor(132, 146, 166)
    add_textbox(slide, "AI Creator Platform · 工程训练营答辩", 0.72, 7.10, 4.2, 0.18, size=7.5, color=color)


def add_card(slide, x, y, w, h, head, body, accent=BLUE, fill=WHITE, num=None):
    rect(slide, x, y, w, h, fill, RGBColor(224, 232, 242), radius=True)
    rect(slide, x, y, 0.08, h, accent, radius=False)
    if num:
        circle(slide, x + 0.28, y + 0.28, 0.34, accent)
        add_textbox(slide, num, x + 0.28, y + 0.36, 0.34, 0.12, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text_x = x + 0.75
    else:
        text_x = x + 0.32
    add_textbox(slide, head, text_x, y + 0.25, w - (text_x - x) - 0.2, min(0.34, max(0.16, h - 0.18)), size=14, color=TEXT, bold=True)
    if body and h > 0.90:
        add_textbox(slide, body, text_x, y + 0.68, w - (text_x - x) - 0.24, h - 0.82, size=9.5, color=MUTED)


def add_stat(slide, x, y, w, h, value, label, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, RGBColor(224, 232, 242), radius=True)
    add_textbox(slide, value, x + 0.22, y + 0.18, w - 0.44, 0.52, size=28, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + 0.20, y + 0.82, w - 0.40, 0.35, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


# 1 封面
s = prs.slides[0]; remove_all_shapes(s); fill_slide(s, NAVY)
rect(s, 8.75, 0, 4.58, 7.5, NAVY_2)
for i, c in enumerate([BLUE, CYAN, RGBColor(80, 140, 255)]):
    rect(s, 8.1 + i * 0.55, 1.0 + i * 0.36, 2.7, 0.34, c, radius=True, transparency=10)
    rect(s, 9.6 + i * 0.38, 3.85 + i * 0.28, 2.2, 0.24, c, radius=True, transparency=25)
add_textbox(s, "AI Creator Platform", 0.88, 1.55, 7.2, 0.70, size=34, color=WHITE, bold=True)
add_textbox(s, "AI 创作者辅助生产与分发平台", 0.92, 2.35, 6.6, 0.45, size=19, color=ICE, bold=True)
add_textbox(s, "大型互联网公司工程训练营 · 项目答辩", 0.94, 3.05, 5.4, 0.28, size=12, color=RGBColor(180, 210, 255))
rect(s, 0.92, 3.62, 2.7, 0.06, CYAN)
add_multi(s, ["从灵感生成到内容审核", "从草稿协作到分发排序", "以工程化闭环支撑 AIGC 内容生产"], 0.92, 4.12, 6.4, 1.25, size=15, color=WHITE)
add_textbox(s, "React · Express · MySQL · Redis · OpenAI-Compatible AI", 0.94, 6.67, 6.9, 0.28, size=10, color=RGBColor(176, 198, 230))

# 2 目录
s = prs.slides[1]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "答辩路线", "以业务问题为入口，展示产品、架构、核心链路与工程化沉淀", idx=2)
sections = [("01", "背景与挑战", "内容团队从选题到发布存在效率、质量与合规三重压力"), ("02", "产品方案", "AI 工作台、Prompt/素材、审核、编辑、分发形成端到端闭环"), ("03", "系统架构", "UmiJS 前端 + Express 服务 + MySQL/Redis + AI Provider"), ("04", "工程实现", "离线草稿、版本回滚、限流日志、Swagger、自动化测试"), ("05", "成果与展望", "沉淀可扩展内容生产底座，规划协作、数据看板与多平台分发")]
for i, (n, h, b) in enumerate(sections):
    y = 1.55 + i * 1.03
    circle(s, 0.98, y + 0.06, 0.46, BLUE if i < 2 else CYAN if i == 2 else NAVY_2)
    add_textbox(s, n, 1.07, y + 0.19, 0.28, 0.12, size=7.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, h, 1.72, y, 2.1, 0.30, size=14, color=TEXT, bold=True)
    add_textbox(s, b, 3.85, y, 7.8, 0.36, size=10.5, color=MUTED)
    line(s, 1.20, y + 0.56, 1.20, y + 0.96, RGBColor(198, 210, 226), 1)
footer(s)

# 3 section
s = prs.slides[2]; remove_all_shapes(s); fill_slide(s, NAVY)
add_textbox(s, "PART 01", 0.92, 1.12, 2.1, 0.32, size=13, color=CYAN, bold=True)
add_textbox(s, "为什么需要 AI 创作者平台", 0.90, 1.74, 8.6, 0.62, size=31, color=WHITE, bold=True)
add_textbox(s, "内容生产正在从单点工具走向平台化闭环：效率提升只是入口，质量、安全与分发才是完整工程命题。", 0.94, 2.72, 8.4, 0.76, size=15.5, color=ICE)
for i, txt in enumerate(["效率", "质量", "安全", "分发"]):
    x = 0.95 + i * 2.15
    rect(s, x, 5.45, 1.55, 0.72, NAVY_2, RGBColor(68, 102, 160), radius=True)
    add_textbox(s, txt, x, 5.67, 1.55, 0.16, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s, True)

# 4 背景痛点
s = prs.slides[3]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "内容生产链路的四类挑战", "大型内容平台场景下，创作不是一次生成，而是一条可控、可追踪、可分发的工程链路", idx=4)
pains = [("创作效率", "选题、成稿、配图、改写依赖人工切换多工具，交付节奏慢。", BLUE), ("质量波动", "爆款标题、内容结构、摘要与封面质量缺少统一评估标准。", CYAN), ("合规风险", "敏感词、低质内容、事实风险需要在发布前被识别和处理。", RED), ("分发割裂", "创作后的排序、曝光、渠道模拟缺少与内容质量联动。", ORANGE)]
for i, (h, b, c) in enumerate(pains):
    x = 0.82 + (i % 2) * 5.85
    y = 1.65 + (i // 2) * 2.15
    add_card(s, x, y, 5.2, 1.55, h, b, c, WHITE, f"0{i+1}")
rect(s, 3.92, 5.93, 5.5, 0.58, NAVY, radius=True)
add_textbox(s, "答辩核心：把 AI 能力落到可运行、可审计、可扩展的内容生产系统中", 4.13, 6.10, 5.1, 0.16, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 5 产品定位
s = prs.slides[4]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "产品定位：AIGC 内容生产闭环平台", "面向个人创作者与中小内容团队，覆盖从灵感到发布的完整路径", idx=5)
rect(s, 0.85, 1.62, 4.7, 4.85, NAVY, radius=True)
add_textbox(s, "一站式 AI 创作中台", 1.18, 2.0, 3.9, 0.38, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, "把生成、审核、编辑、素材、Prompt 与分发能力沉淀为统一工作台。", 1.22, 2.67, 3.75, 0.95, size=13, color=ICE, align=PP_ALIGN.CENTER)
for i, word in enumerate(["生成", "审核", "编辑", "分发"]):
    circle(s, 1.22 + i * 0.95, 4.45, 0.62, [BLUE, CYAN, GREEN, ORANGE][i])
    add_textbox(s, word, 1.22 + i * 0.95, 4.66, 0.62, 0.12, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (h, b) in enumerate([("业务价值", "降低创作门槛，提升内容产能与发布效率"), ("平台价值", "把质量、安全、排序纳入同一套服务体系"), ("工程价值", "端到端链路可观测、可恢复、可演进")]):
    add_card(s, 6.25, 1.74 + i * 1.55, 5.65, 1.18, h, b, [BLUE, CYAN, NAVY_2][i], WHITE)
footer(s)

# 6 用户旅程
s = prs.slides[5]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "用户旅程：从灵感到发布", "围绕创作者真实工作流，把离散功能组织成连续体验", idx=6)
steps = [("灵感输入", "输入选题/关键词"), ("AI 生成", "标题、正文、摘要、配图"), ("风险审核", "规则 + AI 双通道"), ("富文本编辑", "素材、Prompt、草稿"), ("发布分发", "热榜、排序、渠道模拟")]
for i, (h, b) in enumerate(steps):
    x = 0.75 + i * 2.45
    circle(s, x + 0.78, 2.05, 0.76, [BLUE, CYAN, GREEN, ORANGE, NAVY_2][i])
    add_textbox(s, str(i+1), x + 0.78, 2.30, 0.76, 0.16, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s, x, 3.05, 1.95, 1.42, WHITE, RGBColor(224, 232, 242), radius=True)
    add_textbox(s, h, x + 0.18, 3.28, 1.6, 0.22, size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, b, x + 0.18, 3.72, 1.6, 0.32, size=8.8, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        line(s, x + 1.86, 2.43, x + 2.34, 2.43, RGBColor(165, 185, 216), 2)
rect(s, 1.42, 5.45, 10.35, 0.72, NAVY, radius=True)
add_textbox(s, "体验原则：创作者只关注内容决策，平台在后台完成安全、质量和分发策略编排", 1.65, 5.68, 9.9, 0.15, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 7 section
s = prs.slides[6]; remove_all_shapes(s); fill_slide(s, NAVY)
add_textbox(s, "PART 02", 0.92, 1.12, 2.1, 0.32, size=13, color=CYAN, bold=True)
add_textbox(s, "系统架构与核心模块", 0.90, 1.74, 8.0, 0.62, size=31, color=WHITE, bold=True)
add_textbox(s, "前端体验、后端服务、数据存储、AI Provider 与分发能力分层解耦，支撑训练营项目的完整工程复杂度。", 0.94, 2.72, 8.8, 0.76, size=15.5, color=ICE)
for i, txt in enumerate(["Frontend", "API", "Data", "AI"]):
    rect(s, 1.0 + i * 2.1, 5.30, 1.52, 0.66, NAVY_2, RGBColor(68, 102, 160), radius=True)
    add_textbox(s, txt, 1.0 + i * 2.1, 5.52, 1.52, 0.14, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s, True)

# 8 架构
s = prs.slides[7]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "系统架构：分层解耦的内容生产底座", "React/Umi 前端承载创作体验，Express 服务层统一编排业务、数据与 AI 能力", idx=8)
layers = [("前端体验层", "UmiJS 4 · React 18 · Ant Design 5\n创作工作台 / 文章详情 / 用户中心 / 后台审核", BLUE), ("服务与中间件层", "Express 5 · JWT · Joi · Rate Limiter\nController → Service → Model · Swagger · RequestId", CYAN), ("数据与缓存层", "MySQL / Sequelize · Redis\n文章、用户、素材、Prompt、版本、热榜与限流", GREEN), ("AI 能力层", "火山方舟 / ModelScope · OpenAI-Compatible\n生成、审核、评分、配图与本地降级", ORANGE)]
for i, (h, b, c) in enumerate(layers):
    y = 1.55 + i * 1.18
    rect(s, 1.0, y, 11.0, 0.86, WHITE, RGBColor(220, 228, 240), radius=True)
    rect(s, 1.0, y, 0.14, 0.86, c)
    add_textbox(s, h, 1.42, y + 0.18, 2.5, 0.22, size=12.5, color=TEXT, bold=True)
    add_textbox(s, b, 4.05, y + 0.13, 7.4, 0.34, size=9.2, color=MUTED)
    if i < 3:
        line(s, 6.55, y + 0.88, 6.55, y + 1.08, RGBColor(174, 192, 218), 1.6)
footer(s)

# 9 AI 工作台
s = prs.slides[8]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "核心界面：AI 创作工作台", "Creator 页面聚合生成、审核、素材、Prompt、预览与质量分析，是产品体验的主战场", idx=9)
rect(s, 0.88, 1.45, 11.6, 5.55, RGBColor(236, 242, 250), RGBColor(205, 217, 235), radius=True)
rect(s, 1.15, 1.75, 2.45, 4.90, WHITE, RGBColor(220, 228, 240), radius=True)
rect(s, 3.85, 1.75, 4.35, 4.90, WHITE, RGBColor(220, 228, 240), radius=True)
rect(s, 8.45, 1.75, 3.70, 4.90, WHITE, RGBColor(220, 228, 240), radius=True)
add_textbox(s, "输入与配置", 1.38, 2.05, 1.9, 0.22, size=12, color=TEXT, bold=True)
for i, t in enumerate(["选题关键词", "创作风格", "素材引用", "Prompt 模板"]):
    rect(s, 1.38, 2.52 + i * 0.62, 1.82, 0.34, RGBColor(241, 245, 250), radius=True)
    add_textbox(s, t, 1.52, 2.62 + i * 0.62, 1.5, 0.09, size=7.5, color=MUTED)
add_textbox(s, "富文本编辑器", 4.18, 2.05, 2.1, 0.22, size=12, color=TEXT, bold=True)
for i in range(6):
    rect(s, 4.18, 2.55 + i * 0.45, 3.48 - i * 0.13, 0.12, RGBColor(216, 226, 240), radius=True)
rect(s, 4.18, 5.38, 1.38, 0.38, BLUE, radius=True)
add_textbox(s, "生成内容", 4.47, 5.50, 0.8, 0.08, size=7.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, "质量与审核", 8.78, 2.05, 2.1, 0.22, size=12, color=TEXT, bold=True)
for i, (t, c) in enumerate([("风险等级", RED), ("质量评分", GREEN), ("改写建议", ORANGE), ("分发预览", BLUE)]):
    add_card(s, 8.78, 2.48 + i * 0.74, 2.82, 0.52, t, "", c, RGBColor(248, 251, 255))
footer(s)

# 10 功能矩阵
s = prs.slides[9]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "核心功能矩阵", "围绕内容生产闭环划分能力域，每个能力都对应明确的工程模块", idx=10)
items = [("AI 生成", "标题/正文/摘要/配图", BLUE), ("Prompt 管理", "模板、团队、版本", CYAN), ("素材管理", "上传、引用、复用", GREEN), ("审核面板", "风险规则与 AI 复核", RED), ("质量评分", "结构、原创、可读性", ORANGE), ("文章版本", "草稿、发布、回滚", NAVY_2), ("离线恢复", "IndexedDB 自动保存", BLUE), ("热榜分发", "Redis 排序与渠道模拟", CYAN), ("用户反馈", "点赞、收藏、反馈闭环", GREEN)]
for i, (h, b, c) in enumerate(items):
    x = 0.86 + (i % 3) * 4.12
    y = 1.55 + (i // 3) * 1.62
    rect(s, x, y, 3.55, 1.18, WHITE, RGBColor(224, 232, 242), radius=True)
    circle(s, x + 0.28, y + 0.27, 0.34, c)
    add_textbox(s, h, x + 0.78, y + 0.23, 2.35, 0.20, size=11.5, color=TEXT, bold=True)
    add_textbox(s, b, x + 0.78, y + 0.58, 2.35, 0.24, size=8.6, color=MUTED)
footer(s)

# 11 安全质量闭环
s = prs.slides[10]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "AI 安全与质量闭环", "生成结果必须先经过风险识别、质量评分与改写建议，才能进入发布分发链路", idx=11)
center_x, center_y = 6.42, 3.62
circle(s, center_x - 0.85, center_y - 0.85, 1.7, NAVY)
add_textbox(s, "质量\n中枢", center_x - 0.42, center_y - 0.25, 0.84, 0.38, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
nodes = [("规则审核", "敏感词/风险类型", 2.0, 2.1, RED), ("AI 复核", "语义风险与改写", 9.0, 2.1, BLUE), ("质量评分", "结构/原创/可读性", 2.0, 5.0, GREEN), ("人工标注", "评审反馈沉淀", 9.0, 5.0, ORANGE)]
for h, b, x, y, c in nodes:
    rect(s, x, y, 2.45, 0.96, WHITE, RGBColor(224, 232, 242), radius=True)
    rect(s, x, y, 0.10, 0.96, c)
    add_textbox(s, h, x + 0.28, y + 0.18, 1.85, 0.18, size=11.5, color=TEXT, bold=True)
    add_textbox(s, b, x + 0.28, y + 0.52, 1.85, 0.16, size=8.4, color=MUTED)
line(s, 4.45, 2.58, 5.58, 3.18, RGBColor(165, 185, 216), 2)
line(s, 8.98, 2.58, 7.28, 3.18, RGBColor(165, 185, 216), 2)
line(s, 4.45, 5.48, 5.58, 4.23, RGBColor(165, 185, 216), 2)
line(s, 8.98, 5.48, 7.28, 4.23, RGBColor(165, 185, 216), 2)
footer(s)

# 12 排序分发
s = prs.slides[11]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "分发与排序：质量分驱动曝光策略", "内容不是生成后即结束，平台通过热榜、质量分与渠道模拟形成运营闭环", idx=12)
add_card(s, 0.92, 1.66, 3.35, 1.30, "内容输入", "文章热度、发布时间、互动指标、质量评分进入排序模型。", BLUE, WHITE, "01")
add_card(s, 4.93, 1.66, 3.35, 1.30, "Redis 热榜", "缓存排行榜与限流状态，支撑高频查询与快速更新。", CYAN, WHITE, "02")
add_card(s, 8.94, 1.66, 3.35, 1.30, "渠道模拟", "模拟头条/抖音等分发效果，为运营决策提供参考。", GREEN, WHITE, "03")
line(s, 4.28, 2.30, 4.83, 2.30, RGBColor(165, 185, 216), 2)
line(s, 8.29, 2.30, 8.84, 2.30, RGBColor(165, 185, 216), 2)
for i, (v, lab, c) in enumerate([("质量", "内容结构与合规", BLUE), ("热度", "点赞收藏评论", ORANGE), ("时间", "发布新鲜度", GREEN), ("权重", "平台策略参数", CYAN)]):
    add_stat(s, 1.16 + i * 3.0, 4.30, 2.15, 1.30, v, lab, c)
footer(s)

# 13 离线与版本
s = prs.slides[12]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "可靠性设计：离线草稿与版本追踪", "面向真实创作场景，系统要能恢复、可追踪、可回滚", idx=13)
steps = [("30 秒自动保存", "工作台持续保存草稿"), ("IndexedDB 本地持久化", "断网或刷新后可恢复"), ("Dirty 标记与同步", "在线后补偿同步"), ("版本快照", "发布/编辑过程可追溯"), ("一键回滚", "降低误操作成本")]
for i, (h, b) in enumerate(steps):
    x = 0.82 + i * 2.48
    rect(s, x, 2.10, 1.92, 1.32, WHITE, RGBColor(224, 232, 242), radius=True)
    add_textbox(s, f"0{i+1}", x + 0.18, 2.32, 0.45, 0.16, size=10, color=[BLUE, CYAN, GREEN, ORANGE, NAVY_2][i], bold=True)
    add_textbox(s, h, x + 0.18, 2.72, 1.5, 0.22, size=10.3, color=TEXT, bold=True)
    add_textbox(s, b, x + 0.18, 3.08, 1.48, 0.22, size=7.8, color=MUTED)
    if i < 4:
        line(s, x + 1.92, 2.76, x + 2.32, 2.76, RGBColor(165, 185, 216), 1.8)
rect(s, 2.15, 5.25, 9.0, 0.78, NAVY, radius=True)
add_textbox(s, "工程价值：把「写作过程」纳入系统可靠性设计，而不是只保障发布后的数据状态", 2.42, 5.50, 8.45, 0.15, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 14 工程化
s = prs.slides[13]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "工程化能力沉淀", "训练营项目不仅实现功能，更体现后端服务治理、前端状态管理与可验证交付", idx=14)
left = [("鉴权与会话", "JWT access/refresh token，前端 401 自动刷新"), ("限流与日志", "全局 120 req/min + AI 专项限流 + requestId"), ("API 规范", "Swagger 文档、Joi 参数校验、统一响应结构")]
right = [("数据模型", "Sequelize 模型关联与安全 schema upgrade"), ("降级策略", "AI Key 缺失时本地规则与模拟能力保证演示"), ("测试验证", "TypeScript、Jest、Node test、backend check")]
for i, (h, b) in enumerate(left):
    add_card(s, 0.88, 1.62 + i * 1.55, 5.3, 1.16, h, b, [BLUE, CYAN, GREEN][i], WHITE)
for i, (h, b) in enumerate(right):
    add_card(s, 7.05, 1.62 + i * 1.55, 5.3, 1.16, h, b, [ORANGE, NAVY_2, RED][i], WHITE)
line(s, 6.62, 1.68, 6.62, 6.08, RGBColor(202, 214, 232), 1.2)
footer(s)

# 15 数据模型
s = prs.slides[14]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "数据模型：把创作链路落到可追溯对象", "用户、文章、审核、素材、版本和互动数据共同支撑创作闭环", idx=15)
entities = [("User", "账号、角色、认证信息", BLUE), ("Article", "草稿、发布状态、质量评分", CYAN), ("ArticleVersion", "保存、编辑、同步、回滚记录", GREEN), ("AuditLog", "风险类别、合规结论、AI 原始响应", RED), ("Material / Prompt", "素材库、Prompt 模板与团队协作", ORANGE)]
for i, (h, b, c) in enumerate(entities):
    x = 0.86 + (i % 3) * 4.12
    y = 1.62 + (i // 3) * 1.72
    rect(s, x, y, 3.55, 1.28, WHITE, RGBColor(224, 232, 242), radius=True)
    rect(s, x, y, 3.55, 0.10, c, radius=False)
    add_textbox(s, h, x + 0.28, y + 0.30, 2.8, 0.22, size=13, color=TEXT, bold=True)
    add_textbox(s, b, x + 0.28, y + 0.70, 2.95, 0.28, size=8.8, color=MUTED)
line(s, 2.65, 2.98, 4.65, 2.98, RGBColor(166, 187, 216), 1.7)
line(s, 6.80, 2.98, 8.72, 2.98, RGBColor(166, 187, 216), 1.7)
rect(s, 1.20, 5.42, 10.9, 0.64, NAVY, radius=True)
add_textbox(s, "答辩表达重点：不是只做页面，而是把内容生产中的每个状态都沉淀成可审计、可恢复、可扩展的数据对象。", 1.48, 5.63, 10.35, 0.14, size=9.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 16 AI Provider 与降级策略
s = prs.slides[15]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "AI Provider：可切换、可降级、可演示", "文本生成、审核评分、配图与视频入口统一封装在服务层，避免前端直接接触模型密钥", idx=16)
add_card(s, 0.88, 1.62, 3.55, 1.32, "OpenAI-Compatible", "火山方舟与 ModelScope 通过兼容 SDK 接入，Provider 由环境变量切换。", BLUE, WHITE, "01")
add_card(s, 4.72, 1.62, 3.55, 1.32, "多模态入口", "文本、审核、评分走 responses；配图走 Seedream images；视频保留扩展入口。", CYAN, WHITE, "02")
add_card(s, 8.56, 1.62, 3.55, 1.32, "本地兜底", "密钥缺失或外部 API 不可用时，返回规则化结果，保障答辩演示不中断。", GREEN, WHITE, "03")
for i, (label, c) in enumerate([("Prompt", BLUE), ("Service", CYAN), ("Provider", ORANGE), ("Fallback", GREEN)]):
    x = 1.30 + i * 2.85
    circle(s, x, 4.35, 0.74, c)
    add_textbox(s, label, x - 0.08, 4.60, 0.90, 0.13, size=8.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        line(s, x + 0.82, 4.72, x + 2.50, 4.72, RGBColor(166, 187, 216), 2)
rect(s, 1.55, 5.82, 10.0, 0.54, NAVY, radius=True)
add_textbox(s, "工程价值：AI 能力被封装成稳定服务，而不是散落在页面事件里。", 1.78, 6.00, 9.55, 0.12, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 17 现场演示路线
s = prs.slides[16]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "现场演示路线：15 分钟内讲清业务闭环", "从登录到发布，再到审核、热榜与版本回滚，形成一条可操作的演示链路", idx=17)
demo = [("登录进入", "admin / admin123", BLUE), ("AI 创作", "标题、正文、摘要、配图", CYAN), ("安全审核", "风险识别与改写建议", RED), ("保存发布", "草稿、版本、发布状态", GREEN), ("热榜呈现", "质量分 + 互动 + 时间衰减", ORANGE)]
for i, (h, b, c) in enumerate(demo):
    x = 0.76 + i * 2.46
    rect(s, x, 1.86, 1.94, 3.38, WHITE, RGBColor(224, 232, 242), radius=True)
    circle(s, x + 0.67, 2.22, 0.60, c)
    add_textbox(s, f"{i+1}", x + 0.67, 2.42, 0.60, 0.12, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, h, x + 0.18, 3.20, 1.55, 0.22, size=11.2, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, b, x + 0.18, 3.72, 1.55, 0.38, size=8.4, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        line(s, x + 1.94, 2.52, x + 2.34, 2.52, RGBColor(166, 187, 216), 1.8)
rect(s, 1.10, 5.84, 11.0, 0.58, NAVY, radius=True)
add_textbox(s, "讲述策略：每一步都对应一个业务问题和一个工程实现，避免只展示页面。", 1.38, 6.03, 10.45, 0.12, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 18 测试与交付
s = prs.slides[17]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "测试与交付：把可运行性变成答辩证据", "通过类型检查、单元测试、后端语法检查、接口文档和部署说明支撑项目可信度", idx=18)
left = [("TypeScript", "前端类型约束，降低接口字段错配风险"), ("Jest", "API 客户端与模型逻辑单元测试"), ("Node Test", "后端认证边界与服务逻辑测试")]
right = [("Swagger", "接口可浏览、可调试、可交付"), ("Docker Compose", "MySQL + Redis + 后端一键环境"), ("Smoke E2E", "关键链路冒烟验证脚本")]
for i, (h, b) in enumerate(left):
    add_card(s, 0.92, 1.56 + i * 1.50, 5.35, 1.10, h, b, [BLUE, CYAN, GREEN][i], WHITE)
for i, (h, b) in enumerate(right):
    add_card(s, 7.05, 1.56 + i * 1.50, 5.35, 1.10, h, b, [ORANGE, NAVY_2, RED][i], WHITE)
rect(s, 6.55, 1.72, 0.10, 4.46, RGBColor(202, 214, 232), radius=True)
footer(s)

# 19 关键成果与答辩亮点
s = prs.slides[18]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "关键成果：业务闭环与工程闭环同时落地", "项目既能展示产品价值，也能说明系统如何支撑真实创作流程", idx=19)
stats = [("端到端", "生成-审核-编辑-发布闭环", BLUE), ("双通道", "规则审核 + AI 审核", RED), ("双存储", "MySQL 业务数据 + Redis 热榜缓存", CYAN), ("可恢复", "离线草稿 + 版本回滚", GREEN)]
for i, (v, lab, c) in enumerate(stats):
    add_stat(s, 0.95 + i * 3.05, 1.62, 2.35, 1.38, v, lab, c)
add_card(s, 1.12, 4.15, 3.45, 1.28, "产品完整性", "覆盖创作者、内容审核、读者消费和后台治理等关键角色。", BLUE, WHITE)
add_card(s, 4.98, 4.15, 3.45, 1.28, "工程复杂度", "前后端分层、数据库建模、缓存、AI 服务与安全边界均有体现。", CYAN, WHITE)
add_card(s, 8.84, 4.15, 3.45, 1.28, "可演示性", "有降级、有文档、有测试路径，适合答辩现场稳定呈现。", GREEN, WHITE)
footer(s)

# 20 项目边界与优化空间
s = prs.slides[19]; remove_all_shapes(s); fill_slide(s, LIGHT); title(s, "项目边界：清楚知道哪些是真集成，哪些是可扩展点", "答辩时主动说明边界，体现工程判断而不是回避问题", idx=20)
bounds = [("平台分发", "当前为模拟分发 ID，后续可接入头条/西瓜开放平台。", ORANGE), ("审核准确率", "当前有规则与 AI 双通道，后续需要人工标注数据集评测。", RED), ("素材存储", "支持 URL 与本地演示流程，生产环境可扩展 OSS/CDN。", BLUE), ("团队协作", "Prompt 团队和版本已建模，权限细化可继续扩展。", CYAN)]
for i, (h, b, c) in enumerate(bounds):
    x = 0.95 + (i % 2) * 6.0
    y = 1.72 + (i // 2) * 2.18
    add_card(s, x, y, 5.25, 1.48, h, b, c, WHITE, f"0{i+1}")
rect(s, 2.10, 6.02, 9.0, 0.48, NAVY, radius=True)
add_textbox(s, "边界说明不是减分项，而是证明方案具备继续工程化的路线。", 2.42, 6.18, 8.45, 0.10, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# 21 未来规划
s = prs.slides[20]; remove_all_shapes(s); fill_slide(s, NAVY)
add_textbox(s, "未来规划", 0.92, 0.70, 3.0, 0.42, size=24, color=CYAN, bold=True)
add_textbox(s, "从训练营项目走向内容运营中台", 0.92, 1.34, 8.5, 0.58, size=29, color=WHITE, bold=True)
plans = [("多平台分发", "接入头条/西瓜等渠道与投放策略"), ("团队协作", "Prompt 团队、素材权限、审核流转"), ("数据看板", "曝光、互动、转化与运营复盘"), ("模型评测", "审核准确率、生成质量与成本追踪")]
for i, (h, b) in enumerate(plans):
    x = 0.96 + (i % 2) * 5.95
    y = 2.72 + (i // 2) * 1.38
    rect(s, x, y, 5.18, 0.92, NAVY_2, RGBColor(68, 102, 160), radius=True)
    add_textbox(s, h, x + 0.30, y + 0.18, 1.75, 0.18, size=12, color=WHITE, bold=True)
    add_textbox(s, b, x + 2.18, y + 0.18, 2.65, 0.24, size=8.7, color=ICE)
footer(s, True)

# 22 总结 / Q&A
s = prs.slides[21]; remove_all_shapes(s); fill_slide(s, NAVY)
rect(s, 8.70, 0, 4.63, 7.5, NAVY_2)
add_textbox(s, "AI Creator Platform", 0.92, 1.08, 7.0, 0.55, size=30, color=WHITE, bold=True)
add_textbox(s, "用 AI 提升创作效率，用工程化保证内容安全与分发闭环", 0.94, 1.98, 7.2, 0.42, size=17, color=ICE, bold=True)
summary = [("业务闭环", "创作、审核、编辑、发布、消费"), ("工程闭环", "认证、限流、日志、测试、部署"), ("演示闭环", "离线可用、降级可用、路径可复现")]
for i, (h, b) in enumerate(summary):
    rect(s, 1.02, 3.02 + i * 0.82, 6.5, 0.54, NAVY_2, RGBColor(68, 102, 160), radius=True)
    add_textbox(s, h, 1.30, 3.17 + i * 0.82, 1.45, 0.11, size=10.5, color=WHITE, bold=True)
    add_textbox(s, b, 2.90, 3.17 + i * 0.82, 3.85, 0.11, size=9, color=ICE)
add_textbox(s, "谢谢各位老师和同学", 0.92, 6.10, 5.2, 0.38, size=20, color=WHITE, bold=True)
add_textbox(s, "答辩人：熊怀强  |  指导老师：陈佳彬、李超", 0.94, 6.76, 5.8, 0.18, size=8.8, color=RGBColor(176, 198, 230))

prs.save(str(OUT))
print(OUT)
